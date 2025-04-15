import streamlit as st
import requests
from fpdf import FPDF
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import os

# --- Streamlit Setup ---
st.set_page_config(page_title="CAFBrain Generator", layout="centered")
st.title("🧠 CAFBrain: Smart Document Generator")

# --- Initialize Session State Variables ---
if "edited_output" not in st.session_state:
    st.session_state.edited_output = ""
if "edit_history" not in st.session_state:
    st.session_state.edit_history = []
if "redo_stack" not in st.session_state:
    st.session_state.redo_stack = []
if "saved_output" not in st.session_state:
    st.session_state.saved_output = ""
if "sources" not in st.session_state:
    st.session_state.sources = []
if "text_area_version" not in st.session_state:
    st.session_state.text_area_version = 0
if "show_uploader" not in st.session_state:
    st.session_state.show_uploader = False

# --- Export Helpers ---
def generate_pdf(text: str) -> BytesIO:
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("Arial", size=12)
    for line in text.split("\n"):
        pdf.multi_cell(0, 10, line)
    return BytesIO(pdf.output(dest='S').encode('latin-1', errors='ignore'))

def generate_txt(text: str) -> BytesIO:
    return BytesIO(text.encode('utf-8'))

def generate_docx(text: str) -> BytesIO:
    doc = Document()
    for line in text.split("\n"):
        doc.add_paragraph(line)
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def generate_ppt(text: str) -> BytesIO:
    prs = Presentation()
    blocks = text.strip().split("\n\n")
    for i, block in enumerate(blocks):
        lines = [line.strip() for line in block.split("\n") if line.strip()]
        if not lines:
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        # Use title from line 1 or generate default
        title = lines[0] if ":" in lines[0] else f"Slide {i+1}"
        slide.shapes.title.text = title.replace("Slide Title:", "").strip()
        bullets = "\n".join(line.strip("-• ") for line in lines[1:] or lines)
        slide.placeholders[1].text = bullets
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# --- Input Section ---
prompt = st.text_area("📌 Enter your prompt", placeholder="e.g., Create a Canva post about SNAP access")

format_options = {
    "Grant Proposal": "grant",
    "Blog Post": "blog_post",
    "Social Media Post": "social_media_post",
    "Presentation Slides": "presentation",
    "Canva Visual Post": "canva_post",
    "YouTube Video Script": "video_script"
}
selected_label = st.selectbox("📝 Content Type", list(format_options.keys()), index=1)
content_type = format_options[selected_label]

tone = st.selectbox("🎭 Tone", ["formal", "casual", "persuasive", "informative"], index=0)

# --- Document Upload Section ---
# The "+" button reveals the file uploader when clicked.
if st.button("➕ Add Document"):
    st.session_state.show_uploader = True

if st.session_state.show_uploader:
    uploaded_files = st.file_uploader("Upload your documents", accept_multiple_files=True)
    if uploaded_files:
        for file in uploaded_files:
            st.write("Uploaded file:", file.name)

            save_path = f"/Users/sharvari/Downloads/CAFB_Challenge/data/uploads/{file.name}"
            with open(save_path, "wb") as f:
                f.write(file.getbuffer())
            st.success(f"✅ Saved {file.name} to {save_path}")

# Run the update pipeline after upload
    with st.spinner("Processing new documents..."):
        result = os.system("/Users/sharvari/Downloads/CAFB_Challenge/code/update_pipeline.py")
        if result == 0:
            st.success("✅ Documents processed successfully.")
        else:
            st.warning("⚠️ Document processing may have failed. Check console for errors.")


# --- Download formats (before submission)
selected_formats = st.multiselect(
    "💾 Choose file formats to download",
    ["PDF", "DOCX", "TXT", "PPTX"],
    default=["PDF"]
)

# --- Generation Trigger ---
submit = st.button("🚀 Generate Content")

if submit and prompt:
    with st.spinner("Generating content, please wait..."):
        api_url = "http://localhost:8000/generate"
        payload = {
            "query": prompt,
            "top_k": 5,
            "format": content_type,
            "tone": tone
        }
        try:
            response = requests.post(api_url, json=payload)
            if response.status_code == 200:
                data = response.json()
                generated_text = data["answer"]
                sources = data["sources"]
                # Store generated content and sources
                st.session_state.edited_output = generated_text
                st.session_state.edit_history = [generated_text]
                st.session_state.redo_stack = []
                st.session_state.saved_output = ""
                st.session_state.sources = sources
                st.session_state.text_area_version += 1  # update version when new content arrives
            else:
                st.error(f"❌ Backend Error {response.status_code}: {response.text}")
        except Exception as e:
            st.error(f"⚠️ Request failed: {e}")

# --- Editing and Display Section ---
if st.session_state.edited_output:
    st.subheader("📝 Generated Output")
    
    # Create a placeholder for the text area widget
    text_area_placeholder = st.empty()
    
    # Render the text area with a dynamic key based on version
    edited_text = text_area_placeholder.text_area(
        "Edit your content below:",
        key=f"edited_output_{st.session_state.text_area_version}",
        value=st.session_state.edited_output,
        height=300
    )
    
    # Update session state if the user manually edits text
    if edited_text != st.session_state.edited_output:
        st.session_state.edit_history.append(st.session_state.edited_output)
        st.session_state.edited_output = edited_text
        st.session_state.redo_stack.clear()
    
    # Render action buttons in columns
    col1, col2, col3 = st.columns(3)
    if col1.button("↩️ Undo") and st.session_state.edit_history:
        st.session_state.redo_stack.append(st.session_state.edited_output)
        st.session_state.edited_output = st.session_state.edit_history.pop()
        st.session_state.text_area_version += 1  # update version for new key
        text_area_placeholder.empty()  # clear the existing widget
        text_area_placeholder.text_area(
            "Edit your content below:",
            key=f"edited_output_{st.session_state.text_area_version}",
            value=st.session_state.edited_output,
            height=300
        )
    if col2.button("↪️ Redo") and st.session_state.redo_stack:
        st.session_state.edit_history.append(st.session_state.edited_output)
        st.session_state.edited_output = st.session_state.redo_stack.pop()
        st.session_state.text_area_version += 1  # update version for new key
        text_area_placeholder.empty()
        text_area_placeholder.text_area(
            "Edit your content below:",
            key=f"edited_output_{st.session_state.text_area_version}",
            value=st.session_state.edited_output,
            height=300
        )
    if col3.button("💾 Save Edit"):
        st.session_state.saved_output = st.session_state.edited_output
        st.success("✅ Edit saved!")
    
    # Determine the text to download: use saved output if available
    output_to_download = st.session_state.saved_output if st.session_state.saved_output else st.session_state.edited_output

    # --- Download Section ---
    st.subheader("📤 Download Output")
    for file_format in selected_formats:
        file_data = None
        file_name = "cafbrain_output"
        mime_type = "application/octet-stream"
        if file_format == "PDF":
            file_data = generate_pdf(output_to_download)
            file_name += ".pdf"
            mime_type = "application/pdf"
        elif file_format == "DOCX":
            file_data = generate_docx(output_to_download)
            file_name += ".docx"
            mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        elif file_format == "TXT":
            file_data = generate_txt(output_to_download)
            file_name += ".txt"
            mime_type = "text/plain"
        elif file_format == "PPTX":
            file_data = generate_ppt(output_to_download)
            file_name += ".pptx"
            mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        if file_data:
            st.download_button(
                f"⬇️ Download {file_format}",
                data=file_data,
                file_name=file_name,
                mime=mime_type
            )
    
    # --- Sources Section ---
    if st.session_state.sources:
        st.subheader("📚 Sources Used")
        for s in st.session_state.sources:
            st.markdown(f"**{s['title']}**  \n*{s['source']}*  \nScore: {s['score']:.2f}")
            st.caption(s["text"][:400] + "...")
